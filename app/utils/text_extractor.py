import io
import logging
import os
import tempfile
from typing import Dict, Any, List, Optional

# Optional imports for different file types
try:
    import docx

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2

    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False

try:
    import pptx

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)


class TextExtractor:
    """Utility class to extract text from different file types."""

    @staticmethod
    def extract_text_content(file_content: bytes, mime_type: str) -> str:
        """
        Extract text from various file types based on mime_type.

        Args:
            file_content: Binary content of the file
            mime_type: MIME type of the file

        Returns:
            Extracted text content as string
        """
        try:
            if mime_type == 'text/plain':
                return file_content.decode('utf-8', errors='replace')

            elif mime_type == 'application/pdf':
                return TextExtractor._extract_from_pdf(file_content)

            elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                               'application/msword']:
                return TextExtractor._extract_from_docx(file_content)

            elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                               'application/vnd.ms-powerpoint']:
                return TextExtractor._extract_from_pptx(file_content)

            elif mime_type in ['text/html', 'application/xhtml+xml']:
                return TextExtractor._extract_from_html(file_content)

            elif mime_type in ['application/json']:
                return file_content.decode('utf-8', errors='replace')

            elif mime_type in ['text/markdown', 'text/x-markdown']:
                return file_content.decode('utf-8', errors='replace')

            else:
                # Fallback to treating as plain text
                logger.warning(f"Unsupported mime type: {mime_type}, trying to extract as plain text")
                try:
                    return file_content.decode('utf-8', errors='replace')
                except Exception:
                    return "Unable to extract text from this file type."

        except Exception as e:
            logger.error(f"Error extracting text: {str(e)}")
            return f"Error extracting text: {str(e)}"

    @staticmethod
    def _extract_from_pdf(file_content: bytes) -> str:
        """Extract text from a PDF file."""
        if not PYPDF_AVAILABLE:
            logger.warning("PyPDF2 is not installed. Cannot extract text from PDF.")
            return "PDF extraction requires PyPDF2 library. Please install it with: pip install PyPDF2"

        try:
            text = ""
            with io.BytesIO(file_content) as pdf_file:
                reader = PyPDF2.PdfReader(pdf_file)
                for page_num in range(len(reader.pages)):
                    page = reader.pages[page_num]
                    text += page.extract_text() + "\n\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            return f"Error extracting text from PDF: {str(e)}"

    @staticmethod
    def _extract_from_docx(file_content: bytes) -> str:
        """Extract text from a DOCX file."""
        if not DOCX_AVAILABLE:
            logger.warning("python-docx is not installed. Cannot extract text from DOCX.")
            return "DOCX extraction requires python-docx library. Please install it with: pip install python-docx"

        try:
            text = ""
            with io.BytesIO(file_content) as docx_file:
                doc = docx.Document(docx_file)
                for para in doc.paragraphs:
                    text += para.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {str(e)}")
            return f"Error extracting text from DOCX: {str(e)}"

    @staticmethod
    def _extract_from_pptx(file_content: bytes) -> str:
        """Extract text from a PPTX file."""
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx is not installed. Cannot extract text from PPTX.")
            return "PPTX extraction requires python-pptx library. Please install it with: pip install python-pptx"

        try:
            text = ""
            with io.BytesIO(file_content) as pptx_file:
                presentation = pptx.Presentation(pptx_file)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            return f"Error extracting text from PPTX: {str(e)}"

    @staticmethod
    def _extract_from_html(file_content: bytes) -> str:
        """Extract text from an HTML file."""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(file_content, 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.extract()
            # Get text
            text = soup.get_text(separator='\n')
            # Break into lines and remove leading and trailing space on each
            lines = (line.strip() for line in text.splitlines())
            # Break multi-headlines into a line each
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            # Drop blank lines
            text = '\n'.join(chunk for chunk in chunks if chunk)
            return text
        except ImportError:
            logger.warning("BeautifulSoup is not installed. Cannot extract text from HTML properly.")
            return file_content.decode('utf-8', errors='replace')
        except Exception as e:
            logger.error(f"Error extracting text from HTML: {str(e)}")
            return file_content.decode('utf-8', errors='replace')